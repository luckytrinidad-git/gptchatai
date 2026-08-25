import fitz
import json
import io
import numpy as np
from PIL import Image
# import easyocr
from docx import Document
import pandas as pd
import openpyxl
import xlrd
from rapidocr_onnxruntime import RapidOCR
from pathlib import Path
import requests
from gptchatbot.settings import IPFS_SERVER_URL
import re
import yaml
import xlrd
import openpyxl
from pathlib import Path
from email import policy
from email.parser import BytesParser
from docx import Document
from pptx import Presentation
from odf import text, teletype
from odf.opendocument import load
from rapidocr_onnxruntime import RapidOCR
from striprtf.striprtf import rtf_to_text
import extract_msg


# Initialize OCR once
ocr = RapidOCR()


# =========================
# Helper: decode text
# =========================

def decode_text(file_bytes):
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue

    return file_bytes.decode("utf-8", errors="ignore")


# =========================
# Helper: OCR image
# =========================

def ocr_image(img):
    result, _ = ocr(img)

    if not result:
        return ""

    return "\n".join(
        line[1] for line in result
    )


# =========================
# Main extraction function
# =========================

def extract_text(file_name, file_bytes):

    filename = file_name.lower()
    extension = Path(filename).suffix

    # =========================================================
    # TEXT / SOURCE CODE / MARKUP / DATA FILES
    # =========================================================

    text_extensions = {
        ".txt",
        ".md",
        ".csv",
        ".tsv",
        ".log",

        ".html",
        ".htm",
        ".xml",

        ".json",
        ".yaml",
        ".yml",

        ".py",
        ".java",
        ".js",
        ".ts",
        ".jsx",
        ".tsx",
        ".c",
        ".cpp",
        ".h",
        ".cs",
        ".php",
        ".go",
        ".rs",
        ".rb",
        ".swift",
        ".kt",
        ".sql",
        ".sh",
        ".bash",
        ".ps1",
        ".css",
        ".scss",
        ".vue",
    }

    if extension in text_extensions:

        # JSON
        if extension == ".json":
            try:
                data = json.loads(decode_text(file_bytes))
                return json.dumps(data, indent=2, ensure_ascii=False)
            except Exception:
                return decode_text(file_bytes)

        # YAML
        if extension in {".yaml", ".yml"}:
            try:
                data = yaml.safe_load(decode_text(file_bytes))
                return yaml.dump(
                    data,
                    allow_unicode=True,
                    sort_keys=False
                )
            except Exception:
                return decode_text(file_bytes)

        return decode_text(file_bytes)


    # =========================================================
    # RTF
    # =========================================================

    elif extension == ".rtf":

        try:
            return rtf_to_text(
                decode_text(file_bytes)
            )
        except Exception:
            return decode_text(file_bytes)


    # =========================================================
    # PDF
    # VECTOR TEXT + OCR SCANNED PAGES
    # =========================================================

    elif extension == ".pdf":

        output = []

        pdf = fitz.open(
            stream=file_bytes,
            filetype="pdf"
        )

        for page in pdf:

            # ---------------------------------
            # Try normal PDF text extraction
            # ---------------------------------

            page_text = page.get_text("text").strip()

            if page_text:
                output.append(page_text)
                continue

            # ---------------------------------
            # No text -> scanned page
            # ---------------------------------

            pix = page.get_pixmap(
                matrix=fitz.Matrix(1.5, 1.5),
                colorspace=fitz.csRGB,
                alpha=False,
            )

            img = np.frombuffer(
                pix.samples,
                dtype=np.uint8,
            ).reshape(
                pix.height,
                pix.width,
                3,
            )

            page_text = ocr_image(img)

            if page_text:
                output.append(page_text)

        pdf.close()

        return "\n".join(output).strip()


    # =========================================================
    # IMAGE FILES
    # =========================================================

    elif extension in {
        ".jpg",
        ".jpeg",
        ".png",
        ".tif",
        ".tiff",
        ".bmp",
        ".webp",
    }:

        from PIL import Image

        image = Image.open(
            io.BytesIO(file_bytes)
        ).convert("RGB")

        img = np.array(image)

        return ocr_image(img)


    # =========================================================
    # DOCX
    # =========================================================

    elif extension == ".docx":

        doc = Document(
            io.BytesIO(file_bytes)
        )

        output = []

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                output.append(paragraph.text)

        # Tables
        for table in doc.tables:

            for row in table.rows:

                row_text = [
                    cell.text.strip()
                    for cell in row.cells
                ]

                output.append(
                    " | ".join(row_text)
                )

        return "\n".join(output)


    # =========================================================
    # XLSX / XLSM
    # =========================================================

    elif extension in {".xlsx", ".xlsm"}:

        output = []

        wb = openpyxl.load_workbook(
            io.BytesIO(file_bytes),
            data_only=True
        )

        for sheet in wb.sheetnames:

            ws = wb[sheet]

            output.append(
                f"--- Sheet: {sheet} ---"
            )

            for row in ws.iter_rows(
                values_only=True
            ):

                row_text = [
                    str(cell)
                    if cell is not None
                    else ""
                    for cell in row
                ]

                output.append(
                    " | ".join(row_text)
                )

        return "\n".join(output)


    # =========================================================
    # XLS
    # =========================================================

    elif extension == ".xls":

        output = []

        workbook = xlrd.open_workbook(
            file_contents=file_bytes
        )

        for sheet in workbook.sheets():

            output.append(
                f"--- Sheet: {sheet.name} ---"
            )

            for row_idx in range(sheet.nrows):

                row = sheet.row_values(
                    row_idx
                )

                row_text = [
                    str(cell)
                    for cell in row
                ]

                output.append(
                    " | ".join(row_text)
                )

        return "\n".join(output)


    # =========================================================
    # PPTX
    # =========================================================

    elif extension == ".pptx":

        presentation = Presentation(
            io.BytesIO(file_bytes)
        )

        output = []

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1
        ):

            output.append(
                f"--- Slide {slide_number} ---"
            )

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    if shape.text.strip():
                        output.append(
                            shape.text.strip()
                        )

        return "\n".join(output)


    # =========================================================
    # ODT
    # =========================================================

    elif extension == ".odt":

        doc = load(
            io.BytesIO(file_bytes)
        )

        return teletype.extractText(
            doc.text
        )


    # =========================================================
    # ODS
    # =========================================================

    elif extension == ".ods":

        doc = load(
            io.BytesIO(file_bytes)
        )

        output = []

        for table in doc.spreadsheet.getElementsByType(
            __import__("odf.table", fromlist=["Table"]).Table
        ):

            output.append(
                f"--- Sheet: {table.getAttribute('name')} ---"
            )

            for row in table.getElementsByType(
                __import__("odf.table", fromlist=["TableRow"]).TableRow
            ):

                cells = []

                for cell in row.getElementsByType(
                    __import__("odf.table", fromlist=["TableCell"]).TableCell
                ):

                    cells.append(
                        teletype.extractText(cell)
                    )

                output.append(
                    " | ".join(cells)
                )

        return "\n".join(output)


    # =========================================================
    # ODP
    # =========================================================

    elif extension == ".odp":

        doc = load(
            io.BytesIO(file_bytes)
        )

        output = []

        for element in doc.getElementsByType(
            text.P
        ):

            value = teletype.extractText(
                element
            )

            if value.strip():
                output.append(value)

        return "\n".join(output)


    # =========================================================
    # EML
    # =========================================================

    elif extension == ".eml":

        msg = BytesParser(
            policy=policy.default
        ).parsebytes(file_bytes)

        output = []

        if msg["subject"]:
            output.append(
                f"Subject: {msg['subject']}"
            )

        if msg["from"]:
            output.append(
                f"From: {msg['from']}"
            )

        if msg["to"]:
            output.append(
                f"To: {msg['to']}"
            )

        output.append("")

        if msg.is_multipart():

            for part in msg.walk():

                if (
                    part.get_content_type()
                    == "text/plain"
                ):

                    content = part.get_content()

                    if content:
                        output.append(content)

        else:

            if msg.get_content_type() == "text/plain":
                output.append(
                    msg.get_content()
                )

        return "\n".join(output)


    # =========================================================
    # MSG
    # =========================================================

    elif extension == ".msg":

        msg = extract_msg.Message(
            io.BytesIO(file_bytes)
        )

        output = []

        if msg.subject:
            output.append(
                f"Subject: {msg.subject}"
            )

        if msg.sender:
            output.append(
                f"From: {msg.sender}"
            )

        if msg.to:
            output.append(
                f"To: {msg.to}"
            )

        output.append("")

        if msg.body:
            output.append(msg.body)

        msg.close()

        return "\n".join(output)


    # =========================================================
    # Legacy DOC / PPT
    # =========================================================

    elif extension in {".doc", ".ppt"}:

        raise ValueError(
            f"Legacy Office format {extension} requires "
            "LibreOffice conversion before extraction."
        )


    # =========================================================
    # Unknown file
    # =========================================================

    return decode_text(file_bytes)

def chunk_text(text, chunk_size=800, overlap=150):
    chunks = []
    start = 0

    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start = end - overlap

    return chunks

def upload_to_ipfs(uploaded_file):
    """
    uploaded_file:
        Django -> request.FILES["file"]
        FastAPI -> UploadFile.file
        Flask -> request.files["file"]

    Returns:
        {
            "cid": "...",
            "name": "...",
            "size": "..."
        }
    """
    print("Uploading to", f"{IPFS_SERVER_URL}/add")
    print(type(uploaded_file))
    files = {
        "file": (
            uploaded_file.name,
            uploaded_file,
            uploaded_file.content_type
            if hasattr(uploaded_file, "content_type")
            else "application/octet-stream",
        )
    }
    print(files)

    print("Before POST")

    response = requests.post(
        f"{IPFS_SERVER_URL}/add",
        files=files,
        params={"pin": "true"},
        timeout=(10, 300),
    )

    print("After POST")
    print(response.status_code)

    result = response.json()

    cid = result["Hash"]

    return cid

DOCUMENT_MAP = {
    "RAO": "REVENUE ADMINISTRATIVE ORDER",
    "RMO": "REVENUE MEMORANDUM ORDER",
    "RMC": "REVENUE MEMORANDUM CIRCULAR",
    "RR": "REVENUE REGULATIONS",
    "RDAO": "REVENUE DELEGATION AUTHORITY ORDER",
    "RAMO": "REVENUE ADMINISTRATIVE MEMORANDUM ORDER",
    "CMC": "CUSTOMS MEMORANDUM CIRCULAR",
}


def normalize_document_number(number: str) -> str:
    """
    0004-2026
        ↓
    4-2026
    """

    number = number.strip()

    m = re.match(r"^0*(\d+)\s*[-–]\s*(\d{4})$", number)

    if m:
        left = str(int(m.group(1)))
        right = m.group(2).replace(" ", "")
        return f"{left}-{right}"

    return re.sub(r"\s+", "", number)

def build_title_variants(doc_type: str, doc_number: str):
    """
    Example

    RAO
    2-2026

    Generates many searchable title variants.
    """

    doc_number = normalize_document_number(doc_number)

    full_name = DOCUMENT_MAP.get(doc_type)

    variants = set()

    variants.add(f"{doc_type} {doc_number}")
    variants.add(f"{doc_type} NO {doc_number}")
    variants.add(f"{doc_type} NO. {doc_number}")
    variants.add(f"{doc_type} NO.{doc_number}")

    if full_name:

        variants.add(f"{full_name} {doc_number}")
        variants.add(f"{full_name} NO {doc_number}")
        variants.add(f"{full_name} NO. {doc_number}")
        variants.add(f"{full_name} NO.{doc_number}")

    variants.add(doc_number)

    return sorted(variants)

def extract_document_reference(question: str):
    """
    Detects references like

    RAO No. 2-2026

    Revenue Administrative Order No. 2-2026

    RMC 63-2025

    RR No 9-2024

    etc.
    """

    pattern = re.compile(
        r"""
        \b
        (
            RAO|
            RMO|
            RMC|
            RR|
            RDAO|
            RAMO|
            CMC|
            REVENUE\s+ADMINISTRATIVE\s+ORDER|
            REVENUE\s+MEMORANDUM\s+ORDER|
            REVENUE\s+MEMORANDUM\s+CIRCULAR|
            REVENUE\s+REGULATIONS|
            REVENUE\s+DELEGATION\s+AUTHORITY\s+ORDER
        )
        \s*
        (?:NO\.?)?
        \s*
        ([0-9\s\-–]+)
        """,
        re.IGNORECASE | re.VERBOSE,
    )

    match = pattern.search(question)

    if not match:
        return None

    doc_type = match.group(1).upper()

    reverse = {
        "REVENUE ADMINISTRATIVE ORDER": "RAO",
        "REVENUE MEMORANDUM ORDER": "RMO",
        "REVENUE MEMORANDUM CIRCULAR": "RMC",
        "REVENUE REGULATIONS": "RR",
        "REVENUE DELEGATION AUTHORITY ORDER": "RDAO",
    }

    doc_type = reverse.get(doc_type, doc_type)

    doc_number = normalize_document_number(match.group(2))

    return {
        "doc_type": doc_type,
        "doc_number": doc_number,
        "variants": build_title_variants(doc_type, doc_number),
    }